from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# --- THEME CONFIGURATION ---
PRIMARY_COLOR = RGBColor(0, 120, 215)    # Azure Blue
SECONDARY_COLOR = RGBColor(240, 246, 255) # ALMOST White Blue
TEXT_COLOR = RGBColor(50, 50, 50)        # Dark Grey
ACCENT_COLOR = RGBColor(0, 180, 216)     # Cyan Accent

# --- ASSETS PATHS ---
ASSETS_DIR = r"c:/Users/dawse/Desktop/pfa/presentation_assets"
IMG_BG = os.path.join(ASSETS_DIR, "slide_bg_tech_1769983950094.png")
IMG_APP = os.path.join(ASSETS_DIR, "app_concept_phone_1769983968341.png")
IMG_AI = os.path.join(ASSETS_DIR, "ai_neural_diagram_1769983985101.png")
IMG_HW = os.path.join(ASSETS_DIR, "esp32_hardware_art_1769984002383.png")
IMG_HANDS = os.path.join(ASSETS_DIR, "social_inclusion_hands_1769984016858.png")

def create_presentation():
    prs = Presentation()
    
    def add_slide(layout_idx=5): # Blank
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = SECONDARY_COLOR
        return slide

    def add_title(slide, text, subtext=""):
        # Blue Header Bar
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = PRIMARY_COLOR
        shape.line.color.rgb = PRIMARY_COLOR 
        
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1))
        p = tb.text_frame.add_paragraph()
        p.text = text
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.size = Pt(40)
        p.font.bold = True
        
        if subtext:
             p2 = tb.text_frame.add_paragraph()
             p2.text = subtext
             p2.font.color.rgb = RGBColor(220, 220, 220)
             p2.font.size = Pt(18)

    # 1. TITLE SLIDE
    slide = add_slide()
    # Image Background (Full Bleed)
    if os.path.exists(IMG_BG):
        pic = slide.shapes.add_picture(IMG_BG, Inches(0), Inches(0), Inches(10), Inches(7.5))
    
    # Overlay for readability
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = PRIMARY_COLOR
    overlay.fill.transparency = 0.2
    
    # Title Text
    tb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    p = tb.text_frame.add_paragraph()
    p.text = "SignLens"
    p.font.size = Pt(80)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tb.text_frame.add_paragraph()
    p2.text = "L'Intelligence Artificielle au service de l'Inclusion"
    p2.font.size = Pt(32)
    p2.font.color.rgb = RGBColor(240, 240, 240)
    p2.alignment = PP_ALIGN.CENTER

    # 2. CONTEXTE
    slide = add_slide()
    add_title(slide, "1. Contexte & Problématique")
    # Left text, right image (Hands)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(5), Inches(5))
    text = """
    Une fracture invisible :
    
    • 430 millions de personnes souffrent de déficience auditive.
    • La Langue des Signes est leur langue naturelle.
    • 98% des entendants ne la maîtrisent pas.
    
    Résultat : Isolation sociale, difficultés d'accès aux services publics et à l'emploi.
    """
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(20)
    
    if os.path.exists(IMG_HANDS):
        slide.shapes.add_picture(IMG_HANDS, Inches(6), Inches(2.5), Inches(3.5), Inches(3.5))

    # 3. SOLUTION OVERVIEW
    slide = add_slide()
    add_title(slide, "2. La Solution SignLens")
    # Left Image (Phone Mockup), Right Text
    if os.path.exists(IMG_APP):
        slide.shapes.add_picture(IMG_APP, Inches(0.5), Inches(2), Inches(3.5), Inches(5))
        
    tb = slide.shapes.add_textbox(Inches(4.5), Inches(2.5), Inches(5), Inches(4))
    content = """
    Une application mobile révolutionnaire :
    
    ✅ Traduction Instantanée
    (Signes -> Texte/Voix)
    
    ✅ Haute Accessibilité
    (Interface simple, contrastée)
    
    ✅ Hybride & Flexible
    (Caméra Téléphone + Caméra Externe ESP32)
    
    ✅ 100% Hors-Ligne
    (Confidentialité & Rapidité)
    """
    tb.text_frame.text = content
    for paragraph in tb.text_frame.paragraphs:
        paragraph.font.size = Pt(22)
        paragraph.font.color.rgb = TEXT_COLOR

    # 4. ARCHITECTURE LOGIQUE
    slide = add_slide()
    add_title(slide, "3. Architecture Technique")
    # Box Diagram
    # Input
    box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3), Inches(2.5), Inches(2))
    box1.fill.solid()
    box1.fill.fore_color.rgb = RGBColor(230, 230, 230)
    box1.text_frame.text = "Acquisition\nVideo"
    
    arr1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.1), Inches(3.8), Inches(0.8), Inches(0.4))
    arr1.fill.solid()
    arr1.fill.fore_color.rgb = PRIMARY_COLOR
    
    # AI Engine
    box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(2.5), Inches(2.5), Inches(3))
    box2.fill.solid()
    box2.fill.fore_color.rgb = PRIMARY_COLOR
    box2.text_frame.text = "Moteur IA\n(MediaPipe + LSTM)"
    
    arr2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.6), Inches(3.8), Inches(0.8), Inches(0.4))
    arr2.fill.solid()
    arr2.fill.fore_color.rgb = PRIMARY_COLOR

    # Output
    box3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(3), Inches(2.5), Inches(2))
    box3.fill.solid()
    box3.fill.fore_color.rgb = ACCENT_COLOR
    box3.text_frame.text = "Sortie\n(TTS + UI)"

    # 5. ZOOM SUR LE MATERIEL (ESP32)
    slide = add_slide()
    add_title(slide, "4. Hardware : ESP32-CAM")
    
    if os.path.exists(IMG_HW):
        slide.shapes.add_picture(IMG_HW, Inches(5.5), Inches(2), Inches(4), Inches(4))
        
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4.5), Inches(5))
    txt = """
    Pourquoi un module externe ?
    
    • Liberté de mouvement (caméra portable).
    • Coût très faible (< 10$).
    • Streaming MJPEG via WiFi.
    
    Intégration Flutter :
    • Auto-découverte.
    • Latence optimisée.
    """
    tb.text_frame.text = txt
    for p in tb.text_frame.paragraphs:
        p.font.size = Pt(22)

    # 6. ZOOM SUR L'IA (LSTM)
    slide = add_slide()
    add_title(slide, "5. Cerveau IA : Modèle Séquentiel")
    
    if os.path.exists(IMG_AI):
        slide.shapes.add_picture(IMG_AI, Inches(0.5), Inches(2), Inches(5), Inches(4))
        
    tb = slide.shapes.add_textbox(Inches(6), Inches(2), Inches(3.5), Inches(5))
    txt = """
    Architecture LSTM :
    
    • Traite une SEQUENCE temporelle (0.5s).
    • Comprend le mouvement, pas juste la pose.
    • 84 points clés par frame.
    
    Performance :
    • 95% de précision.
    • Inférence < 30ms sur mobile.
    """
    tb.text_frame.text = txt
    for p in tb.text_frame.paragraphs:
        p.font.size = Pt(20)

    # 7. FONCTIONNALITES CLES
    slide = add_slide()
    add_title(slide, "6. Fonctionnalités Clés")
    tb = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
    bullets = [
        "🌐 Mode 100% Hors-Ligne (Embarqué)",
        "🗣️ Synthèse Vocale (Texte -> Parole)",
        "👆 Détection Mains & Doigts (MediaPipe)",
        "🔄 Basculement Caméra Auto (Phone <-> ESP32)",
        "🎨 Interface Moderne (Glassmorphism)"
    ]
    for b in bullets:
        p = tb.text_frame.add_paragraph()
        p.text = b
        p.font.size = Pt(24)
        p.space_before = Pt(24)

    # 8. DEMONSTRATION
    slide = add_slide()
    add_title(slide, "7. Démonstration")
    # Placeholder for video
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2.5), Inches(8), Inches(4))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0, 0, 0)
    box.text_frame.text = "[VIDEO DEMO ICI]"
    box.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # 9. IMPACT & AVENIR
    slide = add_slide()
    add_title(slide, "8. Impact & Perspectives")
    tb = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
    txt = """
    Impact Immédiat :
    • Facilite les échanges simples (Guichets, Commerces).
    • Redonne de l'autonomie.
    
    Futur :
    • Support de phrases complètes.
    • Traduction bi-directionnelle (Voix -> Avatar Signant).
    • Port sur lunettes connectées.
    """
    tb.text_frame.text = txt
    for p in tb.text_frame.paragraphs:
        p.font.size = Pt(22)
        p.space_before = Pt(12)

    # 10. CONCLUSION
    slide = add_slide()
    # Image Background
    if os.path.exists(IMG_HANDS):
        pic = slide.shapes.add_picture(IMG_HANDS, Inches(0), Inches(0), Inches(10), Inches(7.5))
        
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = PRIMARY_COLOR
    overlay.fill.transparency = 0.3
    
    tb = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
    p = tb.text_frame.add_paragraph()
    p.text = "Merci de votre attention"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    prs.save("SignLens_Presentation_Pro.pptx")
    print("Presentation created: SignLens_Presentation_Pro.pptx")

if __name__ == "__main__":
    create_presentation()
